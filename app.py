from flask import Flask,render_template,url_for,request,send_file,redirect,jsonify,session, Response, render_template_string
import time
import spacy
import nltk
import pandas as pd
from spacy_summarization import text_summarizer
from nltk_summarization import nltk_summarizer
from nltk_summarization2 import nltk_summarizer2
from bs4 import BeautifulSoup
from werkzeug.utils import secure_filename
from werkzeug.datastructures import  FileStorage
from heapq import nlargest
import re
import os
nlp = spacy.load('en_core_web_lg')
from spacy.lang.en.stop_words import STOP_WORDS
stopwords = list(STOP_WORDS)
from string import punctuation

import urllib.request
import Text_summarization

import os
import re
import docx
import PyPDF2 
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import json


#app configurations
app=Flask(__name__)
app.config['UPLOAD_FOLDER']= 'c:\\users\\new\\appData\\roaming\\python\\python38\\scripts\\proj1\\f_app\\'

text=''
text_summary=''
pred_summary=''
name=''
global final_summary_nltk2


# import nltk
nltk.download("all")
# exit()

nlp = spacy.load("en_core_web_sm")
timestr = time.strftime("%Y%m%d-%H%M%S")

#here i used Web Scraping Pkg
from bs4 import BeautifulSoup
# instead of this  -- from urllib.request import urlopen i used urllib.request only
import urllib.request


#summarizer
def summariser_spacy(raw_docx):
    raw_text = raw_docx
    docx = nlp(raw_text)
    stopwords = list(STOP_WORDS)
    word_frequencies = {}  
    for word in docx:  
        if word.text not in stopwords:
            if word.text not in word_frequencies.keys():
                word_frequencies[word.text] = 1
            else:
                word_frequencies[word.text] += 1

    maximum_frequncy = max(word_frequencies.values())

    for word in word_frequencies.keys():  
        word_frequencies[word] = (word_frequencies[word]/maximum_frequncy)
    sentence_list = [ sentence for sentence in docx.sents ]

    sentence_scores = {}  
    for sent in sentence_list:  
        for word in sent:
            if word.text.lower() in word_frequencies.keys():
                if len(sent.text.split(' ')) < 100:
                    if sent not in sentence_scores.keys():
                        sentence_scores[sent] = word_frequencies[word.text.lower()]
                    else:
                        sentence_scores[sent] += word_frequencies[word.text.lower()]

    summary_sentences = nlargest(9, sentence_scores, key=sentence_scores.get)
    final_sentences = [ w.text for w in summary_sentences ]
    summary = ' '.join(final_sentences)
    return(summary)

#cleaning function
def clean_text(t1):
    t1=re.sub(r'\[[0-9]*\]',' ',t1)#removing brackets and extra spaces
    t1=re.sub(r'\s+',' ',t1)
    t2=t1
    t2=re.sub('[^a-zA-Z]',' ',t1)#removing special characters and digits
    t2=re.sub(r'\s+',' ',t1)
    return t2

# Fetch Text From Url
def text_from_url(URL):
    article=urllib.request.urlopen(URL)
    parsed_article=BeautifulSoup(article,'html')
    paragraphs=parsed_article.find_all('p')
    article_text=" "
    for p in paragraphs:
        article_text += p.text
    return(str(article_text))

@app.route('/url_text',methods=['GET','POST'])
def url_text():
    if request.method == 'POST':
        raw_url = request.form['raw_url']
        raw_text = text_from_url(raw_url)
        cleaned_text = clean_text(raw_text)
        summary_scraped = summariser_spacy(cleaned_text)
        return render_template('index2.html',summary_scraped=summary_scraped)
   
def readingTime(mytext):
	total_words = len([ token.text for token in nlp(mytext)])
	estimatedTime = total_words/200.0
	return estimatedTime

@app.route('/index2')
def index2():
    return render_template('index2.html')

@app.route('/')
def home():
    return render_template('home.html')
    
@app.route('/index4')
def index4():
    return render_template('index4.html')    

@app.route('/index')
def index():
	return render_template('index.html')

def readfile(flag):
	global text
	global text_summary
	global name
	File=''

	#loads file from 
	if(flag==0):
		File=request.files['filename']
	else:
		File=request.files['file_summary']


	print(File)

	#contains filename
	fn=os.path.basename(File.filename)
	fn=fn.replace(" ", "_")
	print(fn)

	#saving file in specified directory
	File.save(os.path.join(app.config['UPLOAD_FOLDER'],secure_filename(File.filename)))

	#contains path where file is saved
	File=str(os.path.join(app.config['UPLOAD_FOLDER'])+fn)

	#dividing the file location into its filename and extension
	name, ext =os.path.splitext(File)

	#checking whether the file is videofile or not 
	allowed_extensions=[".doc",".docx",".pdf",".rtf",".pptx",".txt"]
	if (ext.lower() in allowed_extensions)==False:
		return "extension_error"

	if(ext.lower()=='.txt'):
		try:
			f = open(File, "r")
			data = f.read()
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.docx' or ext.lower()=='.doc'):
		try:
			doc = docx.Document(File)  # Creating word reader object.
			data = ""
			fullText = []
			for para in doc.paragraphs:
				fullText.append(para.text)
				data = '\n'.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.pptx'):
		try:
			prs = Presentation(File)
			data = ""
			fullText = []
			for slide in prs.slides:
				for shape in slide.shapes:
					if hasattr(shape, "text"):
						fullText.append(shape.text)
				data = '\n'.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.pdf'):
		try:
			pdfFileObj = open(File, mode='rb')    # creating a pdf file object 
			pdfReader = PyPDF2.PdfFileReader(pdfFileObj)    # creating a pdf reader object
			number_of_pages=pdfReader.numPages    #counting number of pages
			data=""
			fullText=[]
			for i in range(0,number_of_pages): 
				pageObj = pdfReader.getPage(i)    # a page object
				fullText.append(pageObj.extractText())   # extracting text from pdf 
			data=' '.join(fullText)
		except Exception:
			return "Something_went_wrong"

	if(ext.lower()=='.rtf'):
		try:
			f = open(File, 'r')
			rtf_text=f.read()
			data = rtf_to_text(rtf_text)
		except Exception:
			return "Something_went_wrong"
			

	# flag=0 -> Original_File
	# flag=1 -> Summary_File
	if(flag==0):
		text=data
	else:
		text_summary=data
	return "successfully_read"


@app.route('/summarize',methods=["GET","POST"])
def summarize():
	global pred_summary

	if request.method=="POST":
		radio_butt=request.form.get("optradio")
		print(radio_butt)
		output=readfile(0)
		
		pred_summary=Text_summarization.get_data(text,'',0)
	return render_template("summarize.html",summary_output=pred_summary)
			
	return render_template('summarize.html')

@app.route('/summerize',methods=['GET','POST'])
def summerizer():
	try:
		start = time.time()
		if request.method == 'POST':
			rawtext = request.json['rawtext']		
			final_summary_nltk = nltk_summarizer2(rawtext)
			len_summary = len(f'{final_summary_nltk}'.split(" "))
			end = time.time()
			final_time = end-start
		return jsonify({'final_summary_nltk': final_summary_nltk, 'final_time': f'{final_time:.3f}', "len_words": f'{len_summary}'})
	except Exception as e:
		print("Error", e)
		return jsonify({'status': f'{e}'})
if __name__ == '__main__':
    app.run(debug=True)
